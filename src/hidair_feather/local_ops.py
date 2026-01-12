import hashlib
import html
import os
import re
from pathlib import Path
from typing import List, Optional

from . import arxiv_ops
from .utils import read_text, safe_filename

try:
    import docx  # type: ignore
except Exception:
    docx = None

try:
    import pptx  # type: ignore
except Exception:
    pptx = None

try:
    from bs4 import BeautifulSoup  # type: ignore
except Exception:
    BeautifulSoup = None

DOCX_AVAILABLE = docx is not None
PPTX_AVAILABLE = pptx is not None
BS4_AVAILABLE = BeautifulSoup is not None

SUPPORTED_EXTS = {".txt", ".md", ".pdf", ".docx", ".pptx", ".html", ".htm"}


def is_supported(path: Path) -> bool:
    return path.suffix.lower() in SUPPORTED_EXTS


def compute_sha1(path: Path) -> str:
    digest = hashlib.sha1()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def build_doc_id(digest: str) -> str:
    return f"local-{digest[:8]}"


def slug_from_path(path: Path, max_len: int = 60) -> str:
    slug = safe_filename(path.stem, max_len=max_len).strip("_")
    return slug or "file"


def extract_text(path: Path) -> Optional[str]:
    ext = path.suffix.lower()
    if ext in {".txt", ".md"}:
        return read_text(path)
    if ext == ".pdf":
        if not arxiv_ops.PYMUPDF_AVAILABLE:
            raise RuntimeError("Missing dependency: pymupdf (pip install pymupdf)")
        return arxiv_ops.pdf_to_text(path)
    if ext == ".docx":
        if not DOCX_AVAILABLE:
            raise RuntimeError("Missing dependency: python-docx (pip install python-docx)")
        doc = docx.Document(path)  # type: ignore[attr-defined]
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    if ext == ".pptx":
        if not PPTX_AVAILABLE:
            raise RuntimeError("Missing dependency: python-pptx (pip install python-pptx)")
        prs = pptx.Presentation(path)  # type: ignore[attr-defined]
        parts: List[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text.strip())
            if texts:
                parts.append(f"Slide {idx}\n" + "\n".join(texts))
        return "\n\n".join(parts)
    if ext in {".html", ".htm"}:
        return html_to_text(read_text(path))
    return None


def html_to_text(html_text: str) -> str:
    if not html_text:
        return ""
    if BS4_AVAILABLE:
        soup = BeautifulSoup(html_text, "html.parser")  # type: ignore[call-arg]
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text(separator="\n")
        return normalize_text(text)

    text = re.sub(r"<script.*?>.*?</script>", "", html_text, flags=re.IGNORECASE | re.DOTALL)
    text = re.sub(r"<style.*?>.*?</style>", "", text, flags=re.IGNORECASE | re.DOTALL)
    text = re.sub(r"<br\s*/?>", "\n", text, flags=re.IGNORECASE)
    text = re.sub(r"</p>", "\n", text, flags=re.IGNORECASE)
    text = re.sub(r"<[^>]+>", "", text)
    text = html.unescape(text)
    return normalize_text(text)


def normalize_text(text: str) -> str:
    text = os.linesep.join(line.strip() for line in text.splitlines())
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()
